from _dummy_thread import exit

import vk
import os
from pptx import Presentation
import datetime
import time
from vk import messages

prs = Presentation("as.pptx")
slide = prs.slides[0]
title = slide.shapes.title

vk.set_access_token('токен юзера')
group = vk.get_group(91865575)

flag = True

timestop = datetime.datetime(2017,4,21, 18,00)

while flag:
    print ("push")
    str = ((timestop - datetime.datetime.now()).seconds / 60).__int__()
    if (str <= 5):
        str = 5
        flag = False
    title.text = str.__str__()
    prs.save("as.pptx")

    os.system("ppt2img.exe as.pptx")

    group.set_cover_photo('as_slide0.png',1590,400)
    print("send"+datetime.datetime.now().strftime('%Y-%m-%d %H:%M'))
    print("без " + str.__str__() + " минут в пути")
    time.sleep(2*60)


